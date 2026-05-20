#!/usr/bin/env python3
"""
build_ppt.py
============
Generates a professional PowerPoint deck that summarises the book.
Run:  python3 scripts/build_ppt.py    ->  build/Mastering-Machine-Learning.pptx

The deck grows as chapters are added: append a (title, [bullets]) entry to
CHAPTER_SUMMARIES and re-run.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- theme -----------------------------------------------------------------
INK    = RGBColor(0x1A, 0x1F, 0x2B)
NAVY   = RGBColor(0x0B, 0x10, 0x26)
PRIMARY= RGBColor(0x4F, 0x46, 0xE5)
SKY    = RGBColor(0x0E, 0xA5, 0xE9)
VIOLET = RGBColor(0x7C, 0x3A, 0xED)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0xCB, 0xD5, 0xE1)

OUT = os.path.join(os.path.dirname(__file__), "..", "build",
                   "Mastering-Machine-Learning.pptx")

prs = Presentation()
prs.slide_width  = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _bar(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def _text(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    r.font.name = "Calibri"
    return tb


def title_slide():
    s = prs.slides.add_slide(BLANK); _bg(s, NAVY)
    _text(s, Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.5),
          "A COMPLETE MASTERY PROGRAM", 14, SKY, bold=True)
    _text(s, Inches(0.8), Inches(1.6), Inches(11.7), Inches(2.0),
          "Mastering Machine Learning", 54, WHITE, bold=True)
    _bar(s, Inches(0.85), Inches(3.95), Inches(4.0), Inches(0.06), SKY)
    _text(s, Inches(0.8), Inches(4.15), Inches(11.7), Inches(0.9),
          "From Beginner to Expert with Theory and Practical Implementation",
          22, GREY)
    _text(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.5),
          "Azhar Hussain", 24, WHITE, bold=True)
    _text(s, Inches(0.8), Inches(6.15), Inches(11.7), Inches(0.6),
          "+92 300 8687258   |   azharhussaincs@gmail.com", 14, GREY)


def section_slide(title, subtitle=""):
    s = prs.slides.add_slide(BLANK); _bg(s, PRIMARY)
    _text(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.6),
          title, 40, WHITE, bold=True)
    if subtitle:
        _text(s, Inches(0.9), Inches(4.2), Inches(11.5), Inches(1.0),
              subtitle, 20, GREY)


def content_slide(title, bullets):
    s = prs.slides.add_slide(BLANK); _bg(s, WHITE)
    _bar(s, 0, 0, SW, Inches(1.15), NAVY)
    _bar(s, 0, Inches(1.15), SW, Inches(0.05), SKY)
    _text(s, Inches(0.7), Inches(0.28), Inches(12), Inches(0.7),
          title, 28, WHITE, bold=True)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.6))
    tf = tb.text_frame; tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        r = p.add_run(); r.text = "•  " + b
        r.font.size = Pt(18); r.font.color.rgb = INK; r.font.name = "Calibri"


# ---- deck content ----------------------------------------------------------
PARTS = [
    "I — Foundations", "II — Maths & Programming Toolkit", "III — Working with Data",
    "IV — Supervised Learning", "V — Unsupervised & Other Paradigms",
    "VI — Deep Learning", "VII — Applied AI Domains",
    "VIII — Production, MLOps & Deployment", "IX — Career, Projects & The Future",
]

# (Section title, chapter title, bullets) — append as chapters are completed.
CHAPTER_SUMMARIES = [
    ("PART I — Foundations", "Chapter 1: Introduction to AI", [
        "AI = making machines do tasks that normally need human intelligence.",
        "AI ⊃ Machine Learning ⊃ Deep Learning (nested fields).",
        "Traditional programming: humans write rules. ML: machines learn rules from data.",
        "Types by capability: Narrow (today), General (not yet), Super (hypothetical).",
        "AI is already around us: maps, recommendations, voice assistants, spam filters.",
        "AI is a tool — its value and its risks depend on how humans use it.",
    ]),
    ("PART I — Foundations", "Chapter 2: Introduction to Machine Learning", [
        "ML (Mitchell): learns from experience E at task T, measured by P — P improves with E.",
        "Modern boom = Big Data + cheap Compute + better Algorithms.",
        "Vocabulary: features (X), labels (y), parameters (learned) vs hyperparameters (you set).",
        "The workflow: define -> collect -> prepare -> split -> choose -> train -> evaluate -> tune -> deploy -> monitor.",
        "Three types: supervised (labelled), unsupervised (no labels), reinforcement (rewards).",
        "Goal = generalization. Beware underfitting (too simple) & overfitting (memorising).",
        "Always evaluate on a separate test set; built a full Iris classifier end-to-end.",
    ]),
    ("PART I — Foundations", "Chapter 3: History of ML", [
        "ML maths predates computers: Bayes (1763), least squares (1805).",
        "1943 neuron, 1950 Turing Test, 1956 'AI' named at Dartmouth, 1957 perceptron.",
        "1969 XOR limitation -> 1st AI winter; 1980s expert systems -> 2nd AI winter.",
        "1986 backpropagation revived multi-layer networks.",
        "2012 AlexNet ignited deep learning; 2017 Transformers enabled modern LLMs.",
        "AI moves in hype/winter cycles — be excited AND skeptical.",
        "Built a perceptron from scratch and saw the XOR limitation first-hand.",
    ]),
    ("PART I — Foundations", "Chapter 4: Types of Machine Learning", [
        "First question of any project: 'Do I have labels?'",
        "Supervised (labels): classification (category) vs regression (number).",
        "Unsupervised (no labels): clustering, dimensionality reduction, anomaly, association.",
        "Semi-supervised = few labels + lots unlabelled; self-supervised = data labels itself (LLMs).",
        "Reinforcement = agent learns a policy from rewards in an environment.",
        "Other splits: batch vs online, instance-based vs model-based.",
        "Ran classification, regression and clustering side by side.",
    ]),
    ("PART II — Maths & Programming Toolkit", "Chapter 5: Mathematics for ML", [
        "Three pillars: linear algebra (store/transform data), calculus (slopes), optimization.",
        "Data = vectors (rows), matrices (datasets), tensors (images/batches).",
        "Dot product powers every prediction: features · weights + bias.",
        "Derivative = slope; gradient = list of slopes, points uphill.",
        "Loss measures 'how wrong'; MSE for regression.",
        "Gradient descent: w <- w - lr * dL/dw, step downhill; learning rate is key.",
        "Coded gradient descent from scratch and watched it learn y = 2x + 1.",
    ]),
    ("PART II — Maths & Programming Toolkit", "Chapter 6: Statistics for ML", [
        "Central tendency: mean (outlier-sensitive), median (robust), mode.",
        "Spread: variance, standard deviation (most used), IQR (robust).",
        "Bayes' theorem: posterior ∝ likelihood × prior; base-rate effect.",
        "Normal curve + 68-95-99.7 rule; z-scores; Central Limit Theorem.",
        "Correlation r ∈ [-1,1] — but correlation is NOT causation.",
        "Inference: samples -> populations; p-values, significance, Type I/II errors.",
        "Computed stats, correlation and a t-test with NumPy/Pandas/SciPy.",
    ]),
    ("PART II — Maths & Programming Toolkit", "Chapter 7: Python for ML", [
        "Python wins for readable syntax, huge free ecosystem, community, notebooks.",
        "Set up: Python 3.10+, pip, a virtual environment per project, Jupyter/VS Code/Colab.",
        "Core: variables, list/tuple/dict/set, if-elif-else, for/while, functions.",
        "Indentation IS syntax; indexing starts at 0.",
        "List comprehensions and lambdas make data code concise.",
        "Classes (self, __init__) underpin every ML library object.",
        "Ecosystem: NumPy, Pandas, Matplotlib/Seaborn, scikit-learn, TF/PyTorch, FastAPI.",
    ]),
    ("PART II — Maths & Programming Toolkit", "Chapter 8: NumPy & Pandas", [
        "NumPy: fast arrays + vectorised maths (10-100x faster than loops).",
        "Key tools: shape/dtype, axis, boolean masks, broadcasting, reshape, dot.",
        "Pandas DataFrame = labelled table; a column is a Series.",
        "Always inspect first: head, info, describe, shape, isnull().sum().",
        "Select/filter (loc vs iloc), create columns, sort_values.",
        "GroupBy = split-apply-combine for per-category stats.",
        "Handle missing values: isnull, fillna, dropna.",
    ]),
    ("PART III — Working with Data", "Chapter 9: Data Analysis Fundamentals", [
        "~80% of ML is understanding & preparing data.",
        "Data types: numerical (discrete/continuous), categorical (nominal/ordinal).",
        "Type dictates valid operations, charts, and encoding.",
        "Sources: CSV, Excel, JSON, SQL, APIs, web, Parquet.",
        "Analytics ladder: descriptive -> diagnostic -> predictive (ML) -> prescriptive.",
        "Analyse univariate -> bivariate -> multivariate; use value_counts, groupby, pivot.",
    ]),
    ("PART III — Working with Data", "Chapter 10: Data Cleaning", [
        "Garbage In, Garbage Out — cleaning sets the ceiling of model quality.",
        "Missing values: detect (isnull), then drop or impute (mean/median/mode/model).",
        "Duplicates: duplicated() + drop_duplicates().",
        "Outliers: detect with IQR (1.5x) or z-score; remove, cap, transform, or keep.",
        "Fix inconsistent text (strip/title/replace) and wrong types (to_numeric/to_datetime).",
        "Compute imputation stats on TRAIN only (avoid leakage); clean deliberately.",
    ]),
    ("PART III — Working with Data", "Chapter 11: Data Preprocessing", [
        "Scaling: normalization (Min-Max -> [0,1]) vs standardization (z-score -> mean0,std1).",
        "Scale for distance/gradient models; NOT for tree-based models.",
        "Encoding: label/ordinal for ordered; one-hot for nominal (no false ordering).",
        "Split train/test FIRST; fit scalers/encoders on train only (avoid leakage).",
        "Use stratify for imbalanced classes.",
        "Pipeline + ColumnTransformer bundle preprocessing + model reproducibly.",
    ]),
    ("PART III — Working with Data", "Chapter 12: Feature Engineering", [
        "Better features beat better algorithms — domain knowledge is the edge.",
        "Create ratios/combinations (BMI, price/m²), differences, counts, flags.",
        "Extract date/time parts (month, weekday, is_weekend).",
        "Binning continuous -> categories; log transform fixes right-skew.",
        "Polynomial & interaction features let linear models capture curves.",
        "High-cardinality categories: frequency/target encoding (not one-hot).",
        "Avoid target leakage; engineer data-dependent features on train only.",
    ]),
    ("PART III — Working with Data", "Chapter 13: Feature Selection", [
        "More features is NOT better — curse of dimensionality.",
        "Benefits: less overfitting, faster, interpretable, often more accurate.",
        "Filter: statistical scores (F-test, correlation, variance) — fast.",
        "Wrapper: search subsets via model (RFE) — accurate, slow.",
        "Embedded: selection during training (Lasso L1, tree importance).",
        "Cross-check methods; select on train/within CV to avoid leakage.",
    ]),
    ("PART III — Working with Data", "Chapter 14: Data Visualization", [
        "Always plot — identical stats can hide different data (Anscombe's Quartet).",
        "Matplotlib (customisable) + Seaborn (statistical, less code).",
        "Six charts: line, bar, histogram, box, scatter, heatmap.",
        "Match chart to the question; title and label everything.",
        "Honest axes: bar charts start at zero; avoid chartjunk.",
        "Correlation heatmaps reveal predictors and redundant features.",
    ]),
    ("PART III — Working with Data", "Chapter 15: Exploratory Data Analysis", [
        "EDA = systematic detective work BEFORE modelling — the Part III capstone.",
        "Workflow: structure -> univariate -> relationships -> quality -> target -> insights.",
        "Titanic: deck 77% missing (drop), age needs imputing, 38% survival baseline.",
        "Sex (74% vs 19%) and class (63%->24%) are powerful survival predictors.",
        "Establish a baseline; pair stats with plots; watch for leakage.",
        "EDA outputs a concrete plan: what to clean, keep, engineer, and which metric.",
    ]),
    ("PART IV — Supervised Learning", "Chapter 16: Supervised Learning Overview", [
        "All supervised algorithms learn f(X) ≈ y from labelled data.",
        "Classification (categories) vs regression (numbers).",
        "Decision boundaries: linear=straight, trees=boxy, KNN=wiggly.",
        "Bias-variance trade-off: too simple=underfit, too complex=overfit; find sweet spot.",
        "No Free Lunch: no algorithm wins everywhere — always compare several.",
        "Bake-off on breast-cancer data: 0.912-0.982; simple linear models tied for best.",
    ]),
    ("PART IV — Supervised Learning", "Chapter 17: Linear Regression", [
        "Predicts a NUMBER as a weighted sum of features: ŷ = w·x + b.",
        "Trains by minimising MSE; via normal equation (exact) or gradient descent.",
        "Metrics: MAE, MSE, RMSE (units), R² (variance explained).",
        "Diabetes data: R²=0.453, RMSE≈53.9 — and coefficients are interpretable.",
        "Assumes linearity; sensitive to outliers & multicollinearity.",
        "Simple, fast, interpretable — the ideal baseline; foundation for later models.",
    ]),
    ("PART IV — Supervised Learning", "Chapter 18: Logistic Regression", [
        "Classification, not regression: linear score + sigmoid -> probability.",
        "Sigmoid squashes z=w·x+b into [0,1]; threshold (0.5) -> class.",
        "Trained with log-loss (cross-entropy), NOT MSE; punishes confident wrong answers.",
        "Outputs calibrated probabilities (predict_proba) and interpretable log-odds.",
        "Multiclass via softmax / one-vs-rest.",
        "Breast-cancer: 0.982 accuracy, precision/recall 0.986; default first classifier.",
    ]),
    ("PART IV — Supervised Learning", "Chapter 19: K-Nearest Neighbors", [
        "Lazy / instance-based: stores data, votes among k nearest neighbours.",
        "Distance metrics: Euclidean (default), Manhattan.",
        "k controls bias-variance: small k overfits, large k underfits; choose via CV.",
        "Feature scaling is CRITICAL — wine accuracy 0.72 -> 0.94 with scaling.",
        "Slow at prediction, memory-heavy, weak in high dimensions.",
        "Great for small, scaled, low-dimensional data and 'similar items' tasks.",
    ]),
    ("PART IV — Supervised Learning", "Chapter 20: Naive Bayes", [
        "Probabilistic classifier from Bayes' theorem + naive independence assumption.",
        "Predicts class maximising P(c) × product of P(xi|c).",
        "Variants: Gaussian (continuous), Multinomial (text counts), Bernoulli (binary).",
        "Laplace smoothing prevents zero-probability wipeouts.",
        "Extremely fast, great on high-dimensional text (spam, sentiment).",
        "Spam demo flagged 'free money now' at 0.947; GaussianNB iris 0.978.",
    ]),
    ("PART IV — Supervised Learning", "Chapter 21: Decision Trees", [
        "Predict by asking yes/no questions root -> leaf; a readable flowchart.",
        "Splits chosen to maximise purity: Gini impurity or entropy/information gain.",
        "Highly interpretable; no scaling needed; gives feature importances.",
        "Overfit easily — control with max_depth, min_samples_leaf, pruning.",
        "Iris: depth=1 underfits (0.667), depth=None overfits (train 1.0, test 0.978).",
        "Unstable alone — the building block of Random Forests & boosting.",
    ]),
    ("PART IV — Supervised Learning", "Chapter 22: Support Vector Machines", [
        "Finds the maximum-margin boundary — widest 'street' between classes.",
        "Support vectors (nearest points) alone define the boundary.",
        "C = soft-margin dial: large C overfits (narrow), small C underfits (wide).",
        "Kernel trick (linear/poly/RBF) draws non-linear boundaries efficiently.",
        "Circles demo: linear 0.41 vs RBF 1.0 — kernels separate the unseparable.",
        "Powerful in high-dim/few-sample; must scale; slow on large data.",
    ]),
    ("PART IV — Supervised Learning", "Chapter 23: Random Forest & Bagging", [
        "Ensembles combine many models: bagging (variance) vs boosting (bias).",
        "Bagging = bootstrap samples + average/vote; cancels random errors.",
        "Random Forest = bagging of trees + random features per split (decorrelates).",
        "Free OOB validation + reliable feature importances; no scaling needed.",
        "Breast-cancer: forest 0.942 > single tree 0.918, OOB 0.967.",
        "Accurate, robust, low-tuning — a premier default for tabular data.",
    ]),
    ("PART IV — Supervised Learning", "Chapter 24: Boosting (AdaBoost, GB, XGBoost)", [
        "Sequential ensemble: each weak learner fixes the previous errors (reduces bias).",
        "AdaBoost reweights misclassified points; Gradient Boosting fits residuals.",
        "XGBoost / LightGBM / CatBoost: fast, regularized — top tabular performers.",
        "Recipe: small learning_rate + many trees + early stopping; shallow weak trees.",
        "CAN overfit (unlike RF) — tune learning_rate & n_estimators together.",
        "Breast-cancer: AdaBoost 0.959, GB 0.947 edged RF 0.942.",
    ]),
    ("PART IV — Supervised Learning", "Chapter 25: Evaluation, Validation & Metrics", [
        "Split train/validation/test; guard the test set; use k-fold cross-validation.",
        "Confusion matrix (TP/FP/TN/FN) is the basis of all classification metrics.",
        "Precision (few false alarms), recall (few misses), F1 (their balance).",
        "ROC/AUC: threshold-independent; 1.0 perfect, 0.5 random.",
        "Accuracy paradox: on imbalanced data accuracy lies — use F1/AUC/recall.",
        "Choose the metric by the COST of each error; never tune on the test set.",
    ]),
    ("PART IV — Supervised Learning", "Chapter 26: Tuning & Regularization", [
        "Tune hyperparameters with grid / random / Bayesian search + cross-validation.",
        "Random search often beats grid per budget (few hyperparameters matter).",
        "Regularization penalises large weights to fight overfitting.",
        "L2 (Ridge) shrinks weights; L1 (Lasso) zeros some (feature selection).",
        "Strength dial: alpha (higher=stronger) or C (lower=stronger).",
        "Grid search found SVM C=10, gamma=0.001 (CV 0.98); Lasso zeroed a feature.",
    ]),
    ("PART V — Unsupervised & Other Paradigms", "Chapter 27: Clustering", [
        "Unsupervised: find structure in UNLABELLED data; clustering groups similar items.",
        "K-Means: assign to k nearest centroids, iterate; choose k via elbow + silhouette.",
        "Hierarchical: dendrogram, no k needed, multi-level (slow on big data).",
        "DBSCAN: density-based — arbitrary shapes, auto-k, detects outliers.",
        "Moons demo: K-Means fails (0.488), DBSCAN finds 2 clusters perfectly.",
        "Scale features; match method to cluster shape; always visualise & interpret.",
    ]),
    ("PART V — Unsupervised & Other Paradigms", "Chapter 28: Dimensionality Reduction", [
        "Compress many features into few: fights curse of dimensionality.",
        "PCA: linear, finds directions of max variance (principal components).",
        "Choose #components via explained variance (digits: 31/64 keep 90%).",
        "t-SNE & UMAP: non-linear, beautiful 2-D VISUALISATIONS (not for models).",
        "Scale before PCA; fit on train, transform both; never feed t-SNE to a classifier.",
        "PCA for compression/preprocessing, t-SNE/UMAP for plots.",
    ]),
    ("PART V — Unsupervised & Other Paradigms", "Chapter 29: Association Rule Learning", [
        "Finds items that go together (market-basket): rules 'if A then B'.",
        "Support = frequency; Confidence = reliability; Lift = vs chance (key metric).",
        "Lift > 1 = positive association; confidence alone can mislead (popular items).",
        "Apriori prunes via 'supersets of rare itemsets are rare'; FP-Growth is faster.",
        "Demo: bread→butter lift 1.11 (real); jam→milk lift 0.89 (illusion).",
        "Rank by lift with min support; association is NOT causation.",
    ]),
    ("PART V — Unsupervised & Other Paradigms", "Chapter 30: Semi-Supervised Learning", [
        "Few labels + lots of unlabelled data — beats the labelling-cost problem.",
        "Works under smoothness / cluster / manifold assumptions.",
        "Self-training: add confident pseudo-labels and retrain.",
        "Label propagation/spreading: labels flow through a similarity graph.",
        "Digits with 131 labels: self-training 0.95 vs supervised-only 0.917.",
        "Self-supervised (auto-generated labels) extends this and powers LLMs.",
    ]),
    ("PART V — Unsupervised & Other Paradigms", "Chapter 31: Reinforcement Learning", [
        "Agent learns by trial, reward & error — no labelled dataset.",
        "Loop: state -> action (policy) -> reward + next state; maximise long-term return.",
        "Discount factor γ values the future; ε-greedy balances explore vs exploit.",
        "Q-learning: Q(s,a) <- Q(s,a) + α[r + γ·maxQ(s',a') - Q(s,a)].",
        "Built a Q-learning agent from scratch; weak exploration broke it first.",
        "Deep RL (DQN, policy gradients) powers Atari, AlphaGo, robotics, RLHF.",
    ]),
    ("PART VI — Deep Learning", "Chapter 32: Neural Networks Foundations", [
        "Neuron: a = φ(w·x + b); a sigmoid neuron is just logistic regression.",
        "Activations add non-linearity: ReLU (hidden default), sigmoid/softmax (output).",
        "MLP = input + hidden + output layers; prediction is the forward pass.",
        "Depth enables hierarchical feature learning (edges->parts->objects).",
        "Built & trained an MLP in PyTorch (0.947 on breast cancer).",
        "Deep learning for unstructured/big data; classic ML often wins on tabular.",
    ]),
    ("PART VI — Deep Learning", "Chapter 33: Training Deep Networks", [
        "Loop: forward -> loss -> backward (backprop) -> update.",
        "Backprop = chain rule run backward; frameworks autodiff via loss.backward().",
        "Optimisers: SGD, momentum, Adam (default); Adam converged ~3x faster than SGD.",
        "Mini-batches + epochs; beware vanishing/exploding gradients (ReLU, batchnorm fix).",
        "Regularize: dropout, batch norm, early stopping, weight decay, augmentation.",
        "Watch train-vs-validation loss; the gap = overfitting -> early stop.",
    ]),
    ("PART VI — Deep Learning", "Chapter 34: Convolutional Neural Networks", [
        "CNNs are built for images: filters convolve to make feature maps.",
        "Local connectivity + parameter sharing -> efficient & spatially aware.",
        "Conv + ReLU + pooling stacks -> hierarchy (edges -> shapes -> objects).",
        "Output size = (W-K+2P)/S+1; max pooling downsamples & adds shift-robustness.",
        "Built a CNN: 0.919 on digits with only 1,898 parameters (efficiency!).",
        "LeNet->AlexNet->VGG->ResNet; use transfer learning for real images.",
    ]),
    ("PART VI — Deep Learning", "Chapter 35: RNNs, LSTM & GRU", [
        "RNNs process sequences step by step, keeping a hidden state (memory).",
        "Plain RNNs forget long-range context (vanishing gradients).",
        "LSTM: cell state + 3 gates (forget/input/output) for long-term memory.",
        "GRU: simpler 2-gate variant, faster; bidirectional adds future context.",
        "Built an LSTM that learned 'is the sum positive' (0.948) — real memory.",
        "Great for time series/streaming; Transformers now dominate NLP.",
    ]),
    ("PART VI — Deep Learning", "Chapter 36: Generative Models", [
        "Generative models learn P(data) to CREATE new data (vs discriminative classify).",
        "Autoencoder: encoder->bottleneck->decoder; compress, denoise, detect anomalies.",
        "AE compressed digits 64->8 (MSE 0.027); VAEs add generation.",
        "GAN: generator vs discriminator (forger vs detective), adversarial training.",
        "GANs are powerful but unstable (mode collapse); diffusion models = modern SOTA.",
        "Dual-use: art & augmentation, but deepfakes/misinfo — build responsibly.",
    ]),
    ("PART VI — Deep Learning", "Chapter 37: Transformers & Attention", [
        "'Attention Is All You Need' (2017) — powers all modern LLMs.",
        "Attention: Query·Key -> softmax weights -> weighted sum of Values.",
        "Self-attention lets every word attend to every other directly (long-range).",
        "Multi-head attention + positional encoding + residual/layernorm blocks.",
        "Won via parallelism + long-range context + scalability.",
        "Families: encoder (BERT, understanding), decoder (GPT, generation), enc-dec (T5).",
    ]),
    ("PART VII — Applied AI Domains", "Chapter 38: Natural Language Processing", [
        "Core challenge: turn text into meaningful numbers.",
        "Preprocess: tokenize, lowercase, stopwords, stemming/lemmatization.",
        "Representations: BoW -> TF-IDF -> word embeddings -> contextual (Transformers).",
        "king - man + woman ≈ queen — embeddings capture meaning as geometry.",
        "Built TF-IDF + Logistic Regression sentiment classifier (correct predictions).",
        "Tasks: classification, sentiment, NER, translation, summarisation, QA.",
    ]),
    ("PART VII — Applied AI Domains", "Chapter 39: Large Language Models", [
        "An LLM = huge Transformer trained to predict the next token.",
        "Training: self-supervised pretraining -> fine-tuning -> RLHF (alignment).",
        "Scale + attention -> emergent abilities; works in tokens within a context window.",
        "Prompting (zero/few-shot, chain-of-thought) steers the same model.",
        "Limitations: hallucination, knowledge cutoff, bias — no true understanding.",
        "Mitigate with RAG (retrieve facts); adapt via prompting -> RAG -> fine-tuning.",
    ]),
    ("PART VII — Applied AI Domains", "Chapter 40: Computer Vision", [
        "Images are grids of pixel numbers; convolution filters detect edges/features.",
        "Tasks: classification, object detection (boxes), segmentation (per-pixel).",
        "CNNs (and Vision Transformers) are the workhorses.",
        "Transfer learning: fine-tune a pretrained model -> strong results, little data.",
        "Data augmentation (flips/rotations/crops) fights overfitting.",
        "Tools: OpenCV, torchvision/Keras; start pretrained, don't train from scratch.",
    ]),
    ("PART VII — Applied AI Domains", "Chapter 41: Recommendation Systems", [
        "Predict which items a user will like out of many (Netflix, Amazon, Spotify).",
        "Content-based: similar items (item features); Collaborative: similar users (ratings).",
        "User-item matrix is sparse; cosine similarity finds taste-mates.",
        "Built user-based CF from scratch (Bob ~ Ann; predicted M5 = 2.69).",
        "Matrix factorization (latent factors) scales (won Netflix Prize).",
        "Challenges: cold start, sparsity, popularity bias, filter bubbles; go hybrid.",
    ]),
    ("PART VII — Applied AI Domains", "Chapter 42: Time Series Forecasting", [
        "Predict future from time-ordered, dependent data — order matters.",
        "Components: trend, seasonality, cyclical, noise; differencing for stationarity.",
        "Methods: moving average, exponential smoothing, ARIMA/SARIMA, ML lags, LSTM.",
        "ML approach: engineer lag/rolling/seasonal/calendar features + a regressor.",
        "NEVER shuffle — split chronologically (train past, test future).",
        "Lag-feature regression MAE 4.13 beat naive baseline 5.30; always beat a baseline.",
    ]),
    ("PART VII — Applied AI Domains", "Chapter 43: Generative AI", [
        "Creates new content: text/code (LLMs), images/video (diffusion), audio.",
        "Diffusion: reverse a noising process — denoise random static into images.",
        "Foundation models: pretrain one giant model, adapt to many tasks.",
        "Temperature dial: low = focused/precise, high = diverse/creative.",
        "Huge opportunities but real risks: hallucination, deepfakes, copyright, bias.",
        "Frontier: agentic AI that takes actions (tools, APIs, multi-step tasks).",
    ]),
    ("PART VIII — Production, MLOps & Deployment", "Chapter 44: Model Deployment", [
        "Deployment makes a trained model available for real predictions.",
        "Serialize the FULL pipeline (preprocessing + model) with joblib; load once at startup.",
        "Serve as a REST API: FastAPI (modern, validated, auto-docs) or Flask (classic).",
        "Streamlit/Gradio for instant interactive demos & dashboards.",
        "Containerise with Docker; pin versions for reproducibility.",
        "Add input validation, logging, monitoring, versioning, rollback.",
    ]),
    ("PART VIII — Production, MLOps & Deployment", "Chapter 45: MLOps", [
        "MLOps = DevOps for ML; production ML is never 'done' — it decays.",
        "ML depends on data + model, not just code; version all three.",
        "Data drift (inputs change) vs concept drift (input->target changes).",
        "KS test detected a shifted batch (p≈0) but passed a similar one (p=0.26).",
        "Components: versioning, experiment tracking, registry, CI/CD/CT, feature store, monitoring.",
        "Monitor from day one; retrain with validation gates; keep a rollback path.",
    ]),
    ("PART VIII — Production, MLOps & Deployment", "Chapter 46: Cloud ML", [
        "Cloud = on-demand GPUs/TPUs, elastic scaling, managed services (pay-as-you-go).",
        "Stack: compute/storage -> managed training -> endpoints -> AutoML -> pre-trained APIs.",
        "Platforms: AWS SageMaker, GCP Vertex AI, Azure ML.",
        "Cost control: spot instances, serverless, shut down idle GPUs, budget alerts.",
        "Watch egress/storage fees and vendor lock-in (containerise).",
        "Use local/edge for privacy, small jobs, or low latency.",
    ]),
    ("PART VIII — Production, MLOps & Deployment", "Chapter 47: Edge AI", [
        "Run models ON the device (phone, camera, car) — low latency, privacy, offline.",
        "Edge devices have limited compute/memory/power -> models must be shrunk.",
        "Quantization (32->8 bit, ~4x smaller, tiny error), pruning, knowledge distillation.",
        "Tools: TensorFlow Lite, ONNX Runtime, Core ML, edge TPUs/NPUs.",
        "Demo: float32 400KB -> int8 100KB (4x) with error 0.0096.",
        "Edge vs cloud is a spectrum; hybrid designs are common.",
    ]),
    ("PART VIII — Production, MLOps & Deployment", "Chapter 48: Responsible AI & Ethics", [
        "Principles: fairness, transparency, privacy, safety, accountability, human oversight.",
        "Bias in -> bias out: models learn & amplify data bias (historical/sampling/proxy).",
        "Measure fairness per subgroup; disparate-impact ratio 0.70 < 0.80 flagged bias.",
        "Accuracy is NOT fairness; explainability (SHAP/LIME) builds trust.",
        "Privacy: consent, differential privacy, federated learning; GDPR, EU AI Act.",
        "Build ethics in from the start; keep humans responsible for big decisions.",
    ]),
    ("PART IX — Career, Projects & The Future", "Chapter 49: Real-World ML Projects", [
        "Projects build skill AND the portfolio that gets you hired.",
        "9-step workflow: define -> data/EDA -> clean -> features -> model -> evaluate -> tune -> deploy -> document.",
        "Worked churn project: pipeline, CV 0.958, test AUC 0.993, saved & deployable.",
        "Catalog of 18 projects (beginner->advanced) mapped to the book's chapters.",
        "Use pipelines (no leakage); FINISH and DEPLOY at least your best projects.",
        "Quality over quantity: 3 polished, documented, deployed projects.",
    ]),
    ("PART IX — Career, Projects & The Future", "Chapter 50: Industry Case Studies", [
        "Same techniques power healthcare, finance, retail, transport, entertainment, etc.",
        "Success pattern: clear problem, good data, right model, deploy, monitor, domain expertise.",
        "Data quality beats model complexity; start simple; tie to a business KPI.",
        "Deployment & monitoring are the hard part; many models never ship or decay.",
        "Most ML projects FAIL: poor framing, bad data/leakage, no deployment, no monitoring.",
        "The algorithm is rarely the bottleneck — data + engineering + framing are.",
    ]),
    ("PART IX — Career, Projects & The Future", "Chapter 51: ML Interview Preparation", [
        "Five areas: concepts, coding, maths/stats, system design, projects/behavioural.",
        "System-design framework: clarify -> problem & metric -> data -> model -> eval -> deploy -> monitor -> ethics.",
        "Question bank distils the book: bias-variance, metrics, regularization, ensembles, drift.",
        "Think out loud, clarify first, start simple, structure open questions.",
        "Know your projects deeply; admit unknowns and reason from fundamentals.",
        "Prepare by reviewing summaries, explaining aloud, and doing mock interviews.",
    ]),
    ("PART IX — Career, Projects & The Future", "Chapter 52: Freelancing with ML", [
        "Earn independently — high demand; you don't need to be an expert.",
        "Services: data analysis/dashboards -> predictive models -> NLP/CV -> LLM/chatbots -> deployment -> consulting.",
        "Find clients: Upwork/Fiverr/Toptal, Kaggle, LinkedIn, niche outreach.",
        "Win with a deployed portfolio + a specialised niche.",
        "Price for VALUE; scope in writing (avoid creep & underpricing).",
        "Deliver usable results, communicate well; repeat clients & referrals sustain you.",
    ]),
    ("PART IX — Career, Projects & The Future", "Chapter 53: Career Guidance & Startups", [
        "Roles: data analyst, data scientist, ML engineer, MLOps, research scientist, data eng, AI PM.",
        "Most roles need fundamentals + portfolio, NOT a PhD (research is the exception).",
        "Choose by strengths; break in via portfolio, networking, contributions, interview prep.",
        "Grow by impact & continuous learning; ML/data roles are highly paid & in demand.",
        "Startups: solve a specific painful problem (ML as enabler), check data, build an MVP.",
        "Foundation models make AI products cheaper to build than ever.",
    ]),
    ("PART IX — Career, Projects & The Future", "Chapter 54: The Future of AI & ML", [
        "Frontiers: multimodal foundation models, AI agents, reasoning, efficient models, AI-for-science.",
        "Emerging: neuro-symbolic, world models, privacy-preserving ML, frontier hardware.",
        "Hardest challenges are non-technical: alignment, safety, interpretability, bias, regulation.",
        "Today's systems are powerful Narrow AI; AGI is hypothetical — keep calibrated judgement.",
        "The #1 durable skill is CONTINUOUS LEARNING grounded in fundamentals.",
        "You are now an ML practitioner — build responsibly, for good. Go build it!",
    ]),
]


def build():
    title_slide()
    section_slide("What This Book Covers",
                  "Nine parts — from zero to production-ready ML engineer")
    content_slide("The 9 Parts of the Journey", PARTS)
    current_section = None
    for section, title, bullets in CHAPTER_SUMMARIES:
        if section != current_section:
            section_slide(section)
            current_section = section
        content_slide(title, bullets)
    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print(f"Saved deck -> {os.path.relpath(OUT)}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
